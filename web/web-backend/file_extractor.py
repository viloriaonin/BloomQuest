import fitz  # PyMuPDF for PDF
from docx import Document
from pptx import Presentation
import openpyxl
import io

def extract_text_from_pdf(file_bytes):
    text = ""
    pdf = fitz.open(stream=file_bytes, filetype="pdf")
    for page in pdf:
        text += page.get_text()
    return text

def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_excel(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text += " ".join([str(c) for c in row if c]) + "\n"
    return text

def extract_text(file_bytes, filename):
    filename = filename.lower()
    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        return extract_text_from_pptx(file_bytes)
    elif filename.endswith(".xlsx") or filename.endswith(".xls"):
        return extract_text_from_excel(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {filename}")